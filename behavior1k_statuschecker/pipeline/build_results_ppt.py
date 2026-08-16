"""Build a results deck: training configs, headline metrics (overall/skill/object
macro recall, base vs. LoRA), per-task breakdowns, example input/output frames, and
one representative frame per demo video (full videos ship in a sibling
<ppt-basename>_videos/ folder - see CLAUDE.md for why they aren't embedded).
Named by date into results/ppt/.
"""
import csv
import datetime
import os
import shutil
from collections import defaultdict

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = "/shared_work/markhsp/Behavior1k-MAPLE/behavior1k_statuschecker"
CSV_DIR = f"{ROOT}/results/csv"
ASSETS_DIR = f"{ROOT}/results/ppt_assets"
PPT_DIR = f"{ROOT}/results/ppt"
os.makedirs(PPT_DIR, exist_ok=True)

# Academic-style palette: restrained, no default-template gradients/rainbow colors.
NAVY = RGBColor(0x1F, 0x33, 0x4A)
SLATE = RGBColor(0x5B, 0x7B, 0x93)
SAGE = RGBColor(0x7A, 0x8B, 0x69)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xB0, 0x2A, 0x2A)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x60, 0x60, 0x60)
TITLE_FONT = "Georgia"
BODY_FONT = "Calibri"
SERIES_COLORS = [NAVY, SLATE, SAGE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_slide_num = [0]


def add_page_number(slide):
    _slide_num[0] += 1
    if _slide_num[0] == 1:
        return  # no page number on the title slide
    box = slide.shapes.add_textbox(Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = str(_slide_num[0])
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.font.name = BODY_FONT


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    add_page_number(slide)
    return slide


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = TITLE_FONT
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = GRAY
        p.font.name = BODY_FONT
    box2 = slide.shapes.add_connector(1, Inches(0.5), Inches(1.15), Inches(12.83), Inches(1.15))
    box2.line.color.rgb = RGBColor(0xC8, 0xC8, 0xC8)
    box2.line.width = Pt(0.75)


def add_bullets(slide, lines, top=1.3, left=0.6, width=12.1, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.font.name = BODY_FONT
        p.space_after = Pt(8)


def add_table(slide, headers, rows, top=1.4, left=0.5, width=12.3, height=5.6, font_size=12):
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = tbl_shape.table
    table.first_row = True
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.name = BODY_FONT
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF1) if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].font.name = BODY_FONT
            cell.text_frame.paragraphs[0].font.color.rgb = DARK
    return table


def add_bar_chart(slide, title, categories, series_dict, top=1.4, left=1.0, width=11.3, height=5.5):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_dict.items():
        chart_data.add_series(name, values)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.name = TITLE_FONT
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = BODY_FONT
    for i, series in enumerate(chart.plots[0].series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = SERIES_COLORS[i % len(SERIES_COLORS)]
    value_axis = chart.value_axis
    value_axis.tick_labels.number_format = "0%"
    value_axis.tick_labels.number_format_is_linked = False
    value_axis.tick_labels.font.size = Pt(11)
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE4, 0xE4, 0xE4)
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(11)
    category_axis.tick_labels.font.name = BODY_FONT


def load_group_csv(name):
    rows = list(csv.DictReader(open(f"{CSV_DIR}/{name}.csv")))
    return rows


def group_avgs(rows):
    """model -> {overall_macro, skill_macro, obj_macro} averaged across tasks"""
    sums = defaultdict(lambda: defaultdict(float))
    counts = defaultdict(int)
    order = []
    for row in rows:
        m = row["model"]
        if m not in order:
            order.append(m)
        counts[m] += 1
        for k in ["overall_macro", "skill_macro", "obj_macro"]:
            sums[m][k] += float(row[k])
    return order, {m: {k: sums[m][k] / counts[m] for k in ["overall_macro", "skill_macro", "obj_macro"]} for m in order}


# ---------------------------------------------------------------- slide 1: title
slide = new_slide()
label = slide.shapes.add_textbox(Inches(1), Inches(2.05), Inches(11.3), Inches(0.4))
lp = label.text_frame.paragraphs[0]
lp.text = "TECHNICAL REPORT"
lp.font.size = Pt(14)
lp.font.color.rgb = SLATE
lp.font.name = BODY_FONT
lp.font.bold = True

box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
tf = box.text_frame
tf.text = "BEHAVIOR-2026 Status Checker"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = TITLE_FONT
tf.paragraphs[0].font.color.rgb = DARK
p = tf.add_paragraph()
p.text = "VLM (LoRA-tuned Qwen3.5-0.8B) per-frame subskill classification: results"
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = GRAY
p.font.name = BODY_FONT
p2 = tf.add_paragraph()
p2.text = datetime.date.today().isoformat()
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0x90, 0x90, 0x90)
p2.font.name = BODY_FONT

# ---------------------------------------------------------------- slide 2: overview
slide = new_slide()
add_title(slide, "What this model does")
add_bullets(slide, [
    "Input: one frame from a fixed head-mounted camera + the task instruction + the task's full label list",
    'Output: the robot\'s current subskill as "<verb> <object>" (e.g. "pick up from radio"), not just a bare verb',
    "Training labels come directly from BEHAVIOR-1K's own skill_annotation ground truth - no manual prompt",
    "engineering or LLM-written CoT required at this scale (classification-only loss)",
    "Base model: Qwen/Qwen3.5-0.8B (native multimodal). Fine-tuning: LoRA (r=16, alpha=32, dropout=0.05,",
    "targeting q/k/v/o/gate/up/down_proj), AdamW lr=1e-4",
    "Evaluated on a fully disjoint held-out episode set (episodes 21-40/task) - never seen during training",
    "Three metrics reported: overall (verb+object combined), skill (verb only), object (object only)",
])

# ---------------------------------------------------------------- slide 3: training configs
slide = new_slide()
add_title(slide, "Training configurations compared")
add_table(slide, ["Run", "Tasks", "Episodes/task (train)", "Epochs", "GPUs", "Examples", "Wall time"], [
    ["100-task", "100", "20 (of up to 200 avail.)", "2", "4 (2 nodes x 2, cross-node DDP)", "~71,700", "~9.5h"],
    ["10-task ablation", "10 (4,10,16,29,37,46,55,64,72,89)", "20", "4", "1 (single-GPU, see gotchas)", "~7,700/epoch", "~4h11m"],
], top=1.5, height=2.0)
add_bullets(slide, [
    "",
    "LoRA config identical across all runs: r=16, alpha=32, dropout=0.05, target_modules=[q/k/v/o/gate/up/down_proj]",
    "Per-class frame cap = 60 (applied after aggregating all episodes for a task, not per-episode)",
    "Checkpoint saved after every epoch - eval below uses epoch1 (\"1 epoch\") / final (\"2 epochs\") for the 100-task",
    "run, and epoch1 (\"2 epochs\") / final (\"4 epochs\") for the 10-task run",
], top=3.8, size=16)

# ---------------------------------------------------------------- slides 4-5: 100-task results
rows_100 = load_group_csv("100task_train1-20ep_eval21-40ep_epoch0-1-2")
order_100, avgs_100 = group_avgs(rows_100)

slide = new_slide()
add_title(slide, "100-task results: base vs. LoRA v3", "Average macro recall across all 100 tasks, held-out eval set")
add_bar_chart(slide, "Avg macro recall by metric", order_100, {
    "overall (verb+object)": [avgs_100[m]["overall_macro"] for m in order_100],
    "skill (verb only)": [avgs_100[m]["skill_macro"] for m in order_100],
    "object (object only)": [avgs_100[m]["obj_macro"] for m in order_100],
})

final_rows_100 = [r for r in rows_100 if r["model"] == order_100[-1]]
best5 = sorted(final_rows_100, key=lambda r: -float(r["overall_macro"]))[:5]
worst5 = sorted(final_rows_100, key=lambda r: float(r["overall_macro"]))[:5]

slide = new_slide()
add_title(slide, "100-task results: best / worst tasks", f"Final model ({order_100[-1]})")
add_table(slide, ["Task", "n_classes", "overall_macro", "skill_macro", "obj_macro"],
          [[r["task_name"], r["n_classes"], r["overall_macro"], r["skill_macro"], r["obj_macro"]] for r in best5],
          top=1.4, height=2.3)
box = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(5), Inches(0.4))
box.text_frame.text = "Best 5 tasks (above); worst 5 (below)"
box.text_frame.paragraphs[0].font.italic = True
add_table(slide, ["Task", "n_classes", "overall_macro", "skill_macro", "obj_macro"],
          [[r["task_name"], r["n_classes"], r["overall_macro"], r["skill_macro"], r["obj_macro"]] for r in worst5],
          top=4.4, height=2.3)

# ---------------------------------------------------------------- slides 6-7: 10-task ablation
rows_10 = load_group_csv("10task_train1-20ep_eval21-40ep_epoch0-2-4")
order_10, avgs_10 = group_avgs(rows_10)

slide = new_slide()
add_title(slide, "10-task ablation: does more training help?", "Same 10 tasks, base vs. 2 epochs vs. 4 epochs")
add_bar_chart(slide, "Avg macro recall by metric", order_10, {
    "overall (verb+object)": [avgs_10[m]["overall_macro"] for m in order_10],
    "skill (verb only)": [avgs_10[m]["skill_macro"] for m in order_10],
    "object (object only)": [avgs_10[m]["obj_macro"] for m in order_10],
})

slide = new_slide()
add_title(slide, "10-task ablation: full per-task breakdown", f"Final model ({order_10[-1]})")
final_rows_10 = sorted([r for r in rows_10 if r["model"] == order_10[-1]], key=lambda r: -float(r["overall_macro"]))
add_table(slide, ["Task", "n_classes", "overall_macro", "skill_macro", "obj_macro"],
          [[r["task_name"], r["n_classes"], r["overall_macro"], r["skill_macro"], r["obj_macro"]] for r in final_rows_10],
          top=1.4, height=5.6, font_size=13)

# ---------------------------------------------------------------- slide 8-9: example in/out
slide = new_slide()
add_title(slide, "Example: correct prediction", "task-0016 (moving_boxes_to_storage), 4-epoch model")
slide.shapes.add_picture(f"{ASSETS_DIR}/example_correct.png", Inches(1.6), Inches(1.1), height=Inches(6.1))

slide = new_slide()
add_title(slide, "Example: wrong prediction", 'task-0029 (clean_up_your_desk) - "move to pen" vs. ground truth "move to pencil"')
slide.shapes.add_picture(f"{ASSETS_DIR}/example_wrong.png", Inches(1.6), Inches(1.1), height=Inches(6.1))
box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.5))
box.text_frame.text = "Correct skill (\"move to\"), wrong object - a near-identical small object confusion, illustrating the skill/object gap seen in the metrics."
box.text_frame.paragraphs[0].font.size = Pt(14)
box.text_frame.paragraphs[0].font.italic = True

# ---------------------------------------------------------------- slides 10-11: demo videos
# Not embedded (see CLAUDE.md: mp4v/FMP4-encoded demo videos aren't reliably
# playable even from genuine PowerPoint) - one representative frame per video
# goes in the deck, the actual mp4 files ship in a sibling <ppt-basename>_videos/
# folder instead.
PPT_BASENAME = f"{datetime.date.today().isoformat()}_statuschecker_results"
VIDEOS_DIR = f"{PPT_DIR}/{PPT_BASENAME}_videos"
os.makedirs(VIDEOS_DIR, exist_ok=True)

for fname, title, subtitle in [
    ("lora_v3_10task4ep_task-0072_cook_a_frozen_pie", "Demo video: cook_a_frozen_pie", "77.1% macro recall, 4-epoch 10-task model"),
    ("lora_v3_10task4ep_task-0016_moving_boxes_to_storage", "Demo video: moving_boxes_to_storage", "76.7% macro recall, 4-epoch 10-task model"),
]:
    slide = new_slide()
    add_title(slide, title, subtitle)
    slide.shapes.add_picture(f"{ASSETS_DIR}/{fname}_poster.png", Inches(2.3), Inches(1.1), height=Inches(5.6))
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4))
    cap = cap_box.text_frame.paragraphs[0]
    cap.text = f"Representative frame only - full video: {PPT_BASENAME}_videos/{fname}.mp4"
    cap.font.size = Pt(13)
    cap.font.italic = True
    cap.font.color.rgb = GRAY
    cap.font.name = BODY_FONT
    shutil.copy(f"{ROOT}/demo/{fname}.mp4", f"{VIDEOS_DIR}/{fname}.mp4")

# ---------------------------------------------------------------- slide 12: next steps
slide = new_slide()
add_title(slide, "Next steps")
add_bullets(slide, [
    "1. Try multi-camera input - only the head camera is used today; wrist views could resolve the",
    "   object-identity ambiguity driving the skill/object gap (see the wrong-prediction example)",
    "2. Try CoT training data - reintroduce reasoning targets (used once in an earlier iteration) at the",
    "   current 100-task scale, especially for the highest-class-count tasks",
    "3. Try adding logical chain constraints - skill sequences aren't unconstrained (e.g. \"move to X\"",
    "   precedes \"pick up from X\"); encoding that structure could clean up sequence-inconsistent predictions",
])

out_path = f"{PPT_DIR}/{PPT_BASENAME}.pptx"
prs.save(out_path)
print(f"wrote {out_path}")
print(f"wrote videos to {VIDEOS_DIR}/")
